#!/usr/bin/env python3
"""
extract_file.py — Universal file content extractor.
Detects file type by extension and extracts readable text content.

Supported formats:
  - PDF      (pdfplumber → PyMuPDF + Tesseract OCR)
  - DOCX     (python-docx)
  - PPTX     (python-pptx)
  - HTML     (BeautifulSoup)
  - CSV/TSV  (pandas or csv module)
  - JSON     (json module, pretty-printed)
  - Archives (list contents only)
  - Plain text (any extension — reads as-is)

Usage:
    python3 extract_file.py <path_to_file>
"""

import sys
import os
import json

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
VENV_PYTHON = os.path.join(SCRIPT_DIR, "..", ".venv", "bin", "python")


# ── Helpers ──────────────────────────────────────────────────────────────────

def get_file_extension(path):
    """Return lowercase file extension without the dot."""
    _, ext = os.path.splitext(path)
    return ext.lower().lstrip(".")


def read_plain_text(path):
    """Read any file as plain text (UTF-8 fallback to latin-1)."""
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(path, "r", encoding="latin-1") as f:
            return f.read()


# ── PDF (pdfplumber → OCR fallback) ─────────────────────────────────────────

def extract_pdf(path):
    """Extract text from PDF. Tries pdfplumber first, falls back to OCR."""
    text = None

    # Try pdfplumber
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    parts.append(f"--- Page {i+1} ---\n{page_text}")
        if parts and any(p.strip() for p in parts):
            text = "\n\n".join(parts)
    except Exception as e:
        print(f"[extract_file] pdfplumber failed: {e}", file=sys.stderr)

    # Fall back to PyMuPDF + Tesseract OCR
    if not text or not text.strip():
        try:
            import pytesseract
            from PIL import Image
            import fitz  # PyMuPDF
            from io import BytesIO

            parts = []
            doc = fitz.open(path)
            for i, page in enumerate(doc):
                mat = fitz.Matrix(2.5, 2.5)
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("png")
                img = Image.open(BytesIO(img_bytes))
                page_text = pytesseract.image_to_string(img)
                parts.append(f"--- Page {i+1} ---\n{page_text}")
            doc.close()
            text = "\n\n".join(parts)
            print("[extract_file] Used OCR fallback (PyMuPDF + Tesseract)", file=sys.stderr)
        except Exception as e:
            print(f"[extract_file] OCR fallback failed: {e}", file=sys.stderr)

    if not text or not text.strip():
        print("ERROR: No text could be extracted from PDF.", file=sys.stderr)
        sys.exit(1)

    return text


# ── DOCX ────────────────────────────────────────────────────────────────────

def extract_docx(path):
    """Extract text from .docx files."""
    try:
        import docx
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs]
        return "\n".join(paragraphs)
    except Exception as e:
        print(f"ERROR extracting DOCX: {e}", file=sys.stderr)
        sys.exit(1)


# ── PPTX ────────────────────────────────────────────────────────────────────

def extract_pptx(path):
    """Extract text from .pptx files."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        print(f"ERROR extracting PPTX: {e}", file=sys.stderr)
        sys.exit(1)


# ── HTML ────────────────────────────────────────────────────────────────────

def extract_html(path):
    """Extract readable text from HTML files."""
    try:
        from bs4 import BeautifulSoup
        with open(path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")

        # Remove script/style elements
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        title = soup.title.string if soup.title else "Untitled"
        body = soup.get_text(separator="\n", strip=True)

        return f"# {title}\n\n{body}"
    except Exception as e:
        print(f"ERROR extracting HTML: {e}", file=sys.stderr)
        sys.exit(1)


# ── CSV / TSV ───────────────────────────────────────────────────────────────

def extract_csv(path, delimiter=","):
    """Extract text summary from CSV/TSV files."""
    try:
        import pandas as pd
        df = pd.read_csv(path, delimiter=delimiter, nrows=None)
        parts = []
        parts.append(f"Shape: {df.shape[0]} rows × {df.shape[1]} columns")
        parts.append(f"\nColumns:\n" + "\n".join(f"  - {col}" for col in df.columns))

        # Summary stats for numeric columns
        numeric_cols = df.select_dtypes(include="number").columns
        if len(numeric_cols) > 0:
            parts.append(f"\nNumeric columns ({len(numeric_cols)}):\n" + ", ".join(numeric_cols))
            stats = df[numeric_cols].describe().to_string()
            parts.append(f"\nSummary Statistics:\n{stats}")

        # Show first rows
        parts.append(f"\nFirst 10 rows:\n{df.head(10).to_string()}")

        return "\n\n".join(parts)
    except ImportError:
        # Fall back to stdlib csv
        import csv
        parts = []
        with open(path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
        parts.append(f"Shape: {len(rows)} rows × {len(rows[0]) if rows else 0} columns")
        if rows:
            parts.append(f"\nHeader: {' | '.join(rows[0])}")
            parts.append(f"\nFirst 10 rows:")
            for row in rows[1:11]:
                parts.append(" | ".join(row))
        return "\n".join(parts)
    except Exception as e:
        return f"ERROR: Could not parse CSV: {e}"


def extract_tsv(path):
    return extract_csv(path, delimiter="\t")


# ── JSON ────────────────────────────────────────────────────────────────────

def extract_json(path):
    """Pretty-print JSON files."""
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"ERROR parsing JSON: {e}", file=sys.stderr)
        sys.exit(1)


# ── Archives ────────────────────────────────────────────────────────────────

def extract_archive(path):
    """List contents of an archive file."""
    try:
        import zipfile
        parts = ["Archive Contents:"]
        with zipfile.ZipFile(path, "r") as z:
            for info in z.infolist():
                parts.append(f"  {info.filename}  ({info.file_size:,} bytes)")
        return "\n".join(parts)
    except zipfile.BadZipFile:
        pass

    # Try tar
    try:
        import tarfile
        parts = ["Archive Contents:"]
        with tarfile.open(path, "r") as tar:
            for member in tar.getmembers():
                parts.append(f"  {member.name}  ({member.size:,} bytes)")
        return "\n".join(parts)
    except Exception:
        pass

    return "ERROR: Unrecognized archive format. Supported: .zip, .tar.gz, .tar.bz2"


# ── Dispatcher ──────────────────────────────────────────────────────────────

# Map extensions to extractor functions
EXTENSIONS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "html": extract_html,
    "htm": extract_html,
    "csv": extract_csv,
    "tsv": extract_tsv,
    "json": extract_json,
    "zip": extract_archive,
    "tar": extract_archive,
    "gz": extract_archive,
    "bz2": extract_archive,
}

BINARY_EXTENSIONS = {
    "exe", "dll", "so", "dylib", "bin", "dat", "db",
    "pyc", "pyo", "pyd", "class", "jar",
    "png", "jpg", "jpeg", "gif", "bmp", "webp", "svg", "ico",
    "mp3", "wav", "ogg", "m4a", "flac", "wma",
    "mp4", "avi", "mkv", "mov", "wmv",
    "ttf", "otf", "woff", "woff2", "eot",
    "psd", "ai", "sketch", "xd", "fig",
}


def main():
    if len(sys.argv) != 2:
        print(f"Usage: {sys.argv[0]} <path_to_file>", file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]
    if not os.path.exists(path):
        print(f"ERROR: File not found: {path}", file=sys.stderr)
        sys.exit(1)

    ext = get_file_extension(path)

    # Check if binary/unprocessable
    if ext in BINARY_EXTENSIONS:
        print(f"NOTE: .{ext} is a binary/image/audio format that cannot be "
              "text-extracted. Use the appropriate tool (vision_analyze for "
              "images, speech-to-text for audio).",
              file=sys.stderr)
        sys.exit(2)

    # Dispatch to the appropriate extractor
    extractor = EXTENSIONS.get(ext)
    if extractor:
        try:
            text = extractor(path)
            print(text)
        except Exception as e:
            print(f"ERROR during extraction: {e}", file=sys.stderr)
            sys.exit(1)
    else:
        # Unknown extension — try as plain text
        try:
            text = read_plain_text(path)
            print(text)
        except Exception as e:
            print(f"ERROR: Could not read file as text: {e}", file=sys.stderr)
            sys.exit(1)


if __name__ == "__main__":
    main()
